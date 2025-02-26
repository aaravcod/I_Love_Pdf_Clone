import sys
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image

def pdf_to_ppt(pdf_path, ppt_path):
    try:
        # Open the PDF file
        doc = fitz.open(pdf_path)
        prs = Presentation()

        for page_num in range(len(doc)):
            page = doc[page_num]
            pix = page.get_pixmap()

            # Convert pixmap to PIL Image
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            # Save to an in-memory buffer
            img_bytes = io.BytesIO()
            img.save(img_bytes, format="PNG")
            img_bytes.seek(0)

            # Add a blank slide
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Add image to slide
            left = top = Inches(0)
            slide.shapes.add_picture(img_bytes, left, top, width=prs.slide_width, height=prs.slide_height)

        # Save PowerPoint file
        prs.save(ppt_path)
        print(f"Converted {pdf_path} to {ppt_path} successfully!")

    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pdf_to_ppt.py <input.pdf> <output.pptx>")
        sys.exit(1)

    input_pdf = sys.argv[1]
    output_ppt = sys.argv[2]

    pdf_to_ppt(input_pdf, output_ppt)
